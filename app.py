import streamlit as st
import streamlit.elements.image as st_image
from PIL import Image, ImageChops, ImageDraw
import io
import re
import zipfile
import base64
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from streamlit_drawable_canvas import st_canvas

# ==========================================
# 0. 页面配置 (必须是第一个 st 命令)
# ==========================================
st.set_page_config(page_title="PDF 瀑布流提取工具", layout="wide", page_icon="📜")

# ==========================================
# 1. 紧急修复补丁 (防止报错)
# ==========================================
if not hasattr(st_image, 'image_to_url'):
    def local_image_to_url(image, width, clamp, channels, output_format, image_id):
        buffered = io.BytesIO()
        if output_format.upper() == "JPEG" and image.mode == "RGBA":
            image = image.convert("RGB")
        image.save(buffered, format=output_format)
        img_str = base64.b64encode(buffered.getvalue()).decode()
        return (f"data:image/{output_format.lower()};base64,{img_str}",)
    st_image.image_to_url = local_image_to_url

# ==========================================
# 2. 核心功能函数
# ==========================================
def sanitize_filename(text):
    text = re.sub(r'\s+', ' ', text).strip()
    return re.sub(r'[\\/*?:"<>|]', "_", text)[:50]

def trim_white_borders(pil_image):
    bg = Image.new(pil_image.mode, pil_image.size, pil_image.getpixel((0,0)))
    diff = ImageChops.difference(pil_image, bg)
    diff = ImageChops.add(diff, diff, 2.0, -100)
    bbox = diff.getbbox()
    if bbox:
        return pil_image.crop(bbox)
    return pil_image

@st.cache_data
def get_page_image(file_content, page_num, zoom=2.0):
    """缓存页面渲染，防止滚动时卡顿"""
    doc = fitz.open(stream=file_content, filetype="pdf")
    page = doc[page_num]
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    doc.close()
    return img

def process_extraction(file_content, page_num, rect_dict, dpi_scale=8.33):
    """处理提取：OCR识别 -> 涂白 -> 裁剪"""
    doc = fitz.open(stream=file_content, filetype="pdf")
    page = doc[page_num]
    
    # 还原坐标 (Canvas 2倍缩放 -> PDF 坐标)
    scale = 0.5 # 因为显示是用2倍缩放的
    r_x = rect_dict["left"] * scale
    r_y = rect_dict["top"] * scale
    r_w = rect_dict["width"] * scale
    r_h = rect_dict["height"] * scale
    rect_pdf = fitz.Rect(r_x, r_y, r_x + r_w, r_y + r_h)
    
    # 1. 提取文字
    text_dict = page.get_text("dict", clip=rect_pdf)
    extracted_text_parts = []
    text_blocks_rects = []
    
    for block in text_dict.get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                text = span["text"].strip()
                if text:
                    extracted_text_parts.append(text)
                    text_blocks_rects.append(span["bbox"])
    
    full_caption = " ".join(extracted_text_parts)
    if not full_caption:
        full_caption = f"Page_{page_num+1}_Image"
        
    # 2. 高清截图
    mat = fitz.Matrix(dpi_scale, dpi_scale)
    pix = page.get_pixmap(matrix=mat, clip=rect_pdf, alpha=False)
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    
    # 3. 涂白文字
    draw = ImageDraw.Draw(img)
    offset_x = rect_pdf.x0
    offset_y = rect_pdf.y0
    
    for bbox in text_blocks_rects:
        x0 = (bbox[0] - offset_x) * dpi_scale
        y0 = (bbox[1] - offset_y) * dpi_scale
        x1 = (bbox[2] - offset_x) * dpi_scale
        y1 = (bbox[3] - offset_y) * dpi_scale
        draw.rectangle([x0-2, y0-2, x1+2, y1+2], fill="white")
        
    # 4. 自动修剪
    final_img = trim_white_borders(img)
    
    out_io = io.BytesIO()
    final_img.save(out_io, format="PNG")
    doc.close()
    
    return out_io.getvalue(), full_caption, final_img.width, final_img.height

# ==========================================
# 3. 界面逻辑
# ==========================================

# 状态初始化
if 'extracted_list' not in st.session_state:
    st.session_state.extracted_list = []

# 小红书笔记链接
XHS_LINK = "https://www.xiaohongshu.com/explore/696f27e7000000000a03ee45?xsec_token=ABft3QO37w_LDTt8J5zePSaog2TSYY1qVxGckdEZeuUpc=&xsec_source=pc_user"

# 统一读取文件内容，优化内存，避免重复 read()
bytes_data = None
total_pages = 0

# --- 侧边栏 ---
with st.sidebar:
    st.header("1. 导入 PDF")
    uploaded_file = st.file_uploader("文件上传", type="pdf")
    
    display_range = None
    if uploaded_file:
        bytes_data = uploaded_file.getvalue()
        doc_temp = fitz.open(stream=bytes_data, filetype="pdf")
        total_pages = len(doc_temp)
        doc_temp.close() # 及时释放内存
        
        # 侧边栏小红书提醒
        st.markdown(f"🌟 **支持作者**：[去小红书留言口令]({XHS_LINK})")
        
        if total_pages > 5:
            st.info(f"文档共 {total_pages} 页")
            display_range = st.slider("显示页码范围 (防止卡顿)", 1, total_pages, (1, min(10, total_pages)))
    
    st.divider()
    st.header("2. 导出结果")
    st.write(f"已提取图片: **{len(st.session_state.extracted_list)}** 张")
    
    # 预览小图
    if st.session_state.extracted_list:
        with st.expander("查看已提取列表"):
            for idx, item in enumerate(st.session_state.extracted_list):
                col_del, col_txt = st.columns([1, 4])
                with col_txt:
                    st.caption(f"{idx+1}. {item['name']}")
    
    if st.button("🗑️ 清空所有"):
        st.session_state.extracted_list = []
        st.rerun()

    # 导出按钮
    if st.session_state.extracted_list:
        c1, c2 = st.columns(2)
        
        # PPTX
        prs = Presentation()
        prs.slide_width = Inches(7.5); prs.slide_height = Inches(10)
        
        for item in st.session_state.extracted_list:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            pw, ph = prs.slide_width, prs.slide_height
            margin = Inches(0.5)
            
            img_io = io.BytesIO(item["bytes"])
            
            # 布局计算
            max_h = ph - Inches(1.5)
            max_w = pw - margin * 2
            ratio = item["w"] / item["h"]
            target_w = max_w
            target_h = target_w / ratio
            if target_h > max_h:
                target_h = max_h
                target_w = target_h * ratio
            
            left = (pw - target_w) / 2
            top = Inches(0.5)
            
            slide.shapes.add_picture(img_io, left, top, width=target_w, height=target_h)
            
            tb = slide.shapes.add_textbox(margin, top + target_h + Inches(0.1), pw - margin*2, Inches(1))
            p = tb.text_frame.add_paragraph()
            p.text = item["name"]
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = "Microsoft YaHei"
            
        ppt_io = io.BytesIO()
        prs.save(ppt_io); ppt_io.seek(0)
        c1.download_button("📥 PPTX", ppt_io, "export.pptx")
        
        # ZIP
        zip_io = io.BytesIO()
        with zipfile.ZipFile(zip_io, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, item in enumerate(st.session_state.extracted_list):
                zf.writestr(f"{i+1}_{item['name']}.png", item["bytes"])
        zip_io.seek(0)
        c2.download_button("📦 ZIP", zip_io, "images.zip")

# --- 主界面：瀑布流显示 ---
st.title("📜 浏览模式提取工具")
st.info("操作方式：像看书一样往下滑，看到想提取的图，直接**画框**，然后点下方的**⚡提取**按钮。")

if uploaded_file and bytes_data:
    # 强制留言提醒区块
    st.success("🎉 **文件读取成功！**")
    st.warning(f"""
    **📢 温馨提示：**  
    如果您觉得这个工具好用，**请务必前往原笔记评论区留言口令**，您的支持是我更新的最大动力！❤️  
    👉 [点击这里一键跳转至小红书原笔记留言]({XHS_LINK})
    """)
    
    # 确定显示范围
    start_p = 0
    end_p = total_pages
    if display_range:
        start_p = display_range[0] - 1
        end_p = display_range[1]
    
    # === 循环渲染每一页 ===
    for p_idx in range(start_p, end_p):
        st.divider()
        st.markdown(f"### 第 {p_idx + 1} 页")
        
        # 1. 获取背景图 (带缓存，速度快)
        bg_image = get_page_image(bytes_data, p_idx)
        
        # 2. 创建画布
        # key 必须唯一，使用页码区分
        canvas_result = st_canvas(
            fill_color="rgba(255, 0, 0, 0.1)",
            stroke_width=2,
            stroke_color="#FF0000",
            background_image=bg_image,
            update_streamlit=True,
            height=bg_image.height,
            width=bg_image.width,
            drawing_mode="rect",
            key=f"canvas_page_{p_idx}", 
            display_toolbar=True,
        )
        
        # 3. 提取按钮 (跟随在每一页下面)
        if canvas_result.json_data and canvas_result.json_data["objects"]:
            last_obj = canvas_result.json_data["objects"][-1]
            
            col_btn, col_msg = st.columns([1, 4])
            with col_btn:
                if st.button(f"⚡ 提取第 {p_idx+1} 页选中区域", key=f"btn_{p_idx}", type="primary"):
                    try:
                        img_bytes, img_name, w, h = process_extraction(bytes_data, p_idx, last_obj)
                        
                        st.session_state.extracted_list.append({
                            "bytes": img_bytes,
                            "name": sanitize_filename(img_name),
                            "page": p_idx + 1,
                            "w": w, "h": h
                        })
                        st.toast(f"✅ 已成功提取: {img_name}")
                        # 强制刷新侧边栏
                        st.rerun()
                    except Exception as e:
                        st.error(f"提取出错: {e}")
            with col_msg:
                st.caption("✅ 已选中区域，点击左侧按钮提取")

else:
    st.warning("👈 请在左侧上传 PDF 文件即可开始。")
